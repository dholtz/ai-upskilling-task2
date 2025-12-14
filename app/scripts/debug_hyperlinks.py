"""Debug script to understand how hyperlinks are stored in PowerPoint"""
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.table import Table

# Get file path from command line argument or use default test path
if len(sys.argv) > 1:
    pptx_path = sys.argv[1]
else:
    # Default test path (can be overridden via command line)
    pptx_path = "/tmp/presentation.pptx"

if not os.path.exists(pptx_path):
    print(f"Error: File not found: {pptx_path}")
    print("Usage: python debug_hyperlinks.py [path_to_pptx_file]")
    sys.exit(1)

if not os.path.exists(pptx_path):
    print(f"File not found: {pptx_path}")
    sys.exit(1)

prs = Presentation(pptx_path)

print(f"Total slides: {len(prs.slides)}\n")

for slide_num, slide in enumerate(prs.slides, 1):
    print(f"\n=== Slide {slide_num} ===")
    
    for shape_num, shape in enumerate(slide.shapes, 1):
        print(f"\nShape {shape_num}: {type(shape).__name__}")
        
        # Check if it's a table
        if isinstance(shape, Table):
            print("  Is a TABLE")
            for row_idx, row in enumerate(shape.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_text = cell.text.strip()
                    if cell_text:
                        print(f"    Cell [{row_idx},{col_idx}]: {cell_text[:60]}")
                        # Check for hyperlinks in table cell
                        if hasattr(cell, "text_frame"):
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if hasattr(run, "hyperlink") and run.hyperlink:
                                        hyperlink = run.hyperlink
                                        url = getattr(hyperlink, 'address', None)
                                        print(f"      Has hyperlink, address: {url}")
                                        if not url and hasattr(hyperlink, '_hlinkClick'):
                                            hlink_click = hyperlink._hlinkClick
                                            if hlink_click and hasattr(hlink_click, 'rId'):
                                                print(f"      Has rId: {hlink_click.rId}")
        
        if hasattr(shape, "text_frame"):
            print("  Has text_frame")
            for para_num, paragraph in enumerate(shape.text_frame.paragraphs, 1):
                para_text = "".join(run.text for run in paragraph.runs if run.text)
                if para_text.strip():
                    print(f"    Paragraph {para_num}: {para_text[:80]}...")
                
                # Check runs for hyperlinks
                for run_num, run in enumerate(paragraph.runs, 1):
                    if hasattr(run, "hyperlink") and run.hyperlink:
                        hyperlink = run.hyperlink
                        print(f"      Run {run_num} has hyperlink:")
                        print(f"        Run text: {run.text[:50] if run.text else 'None'}")
                        print(f"        Has address: {hasattr(hyperlink, 'address')}")
                        if hasattr(hyperlink, 'address'):
                            print(f"        Address: {hyperlink.address}")
                        print(f"        Has rId: {hasattr(hyperlink, 'rId')}")
                        if hasattr(hyperlink, 'rId'):
                            print(f"        rId: {hyperlink.rId}")
                            # Try to resolve rId through slide part
                            try:
                                if hasattr(slide, 'part') and hasattr(slide.part, 'related_parts'):
                                    related_parts = slide.part.related_parts
                                    if hyperlink.rId in related_parts:
                                        rel_part = related_parts[hyperlink.rId]
                                        print(f"        Found related part: {rel_part}")
                                        print(f"        Related part type: {type(rel_part)}")
                                        if hasattr(rel_part, 'target_ref'):
                                            print(f"        Related part target_ref: {rel_part.target_ref}")
                                        if hasattr(rel_part, 'target_uri'):
                                            print(f"        Related part target_uri: {rel_part.target_uri}")
                            except Exception as e:
                                print(f"        Error resolving rId: {e}")
                        print(f"        Hyperlink type: {type(hyperlink)}")
                        # Check hyperlink's internal attributes and XML
                        if hasattr(hyperlink, '_hlinkClick'):
                            hlink_click = hyperlink._hlinkClick
                            print(f"        _hlinkClick: {hlink_click}")
                            if hlink_click and hasattr(hlink_click, 'rId'):
                                print(f"        _hlinkClick.rId: {hlink_click.rId}")
                            if hlink_click and hasattr(hlink_click, 'action'):
                                print(f"        _hlinkClick.action: {hlink_click.action}")
                        
                        # Check the underlying XML element
                        if hasattr(hyperlink, '_element'):
                            element = hyperlink._element
                            print(f"        XML element: {element}")
                            print(f"        XML element tag: {element.tag}")
                            print(f"        XML element attrib: {element.attrib}")
                            # Check for rId in XML
                            if 'rId' in element.attrib:
                                rId = element.attrib['rId']
                                print(f"        XML rId: {rId}")
                                # Try to resolve through slide part
                                try:
                                    if hasattr(slide, 'part') and hasattr(slide.part, 'related_parts'):
                                        related_parts = slide.part.related_parts
                                        if rId in related_parts:
                                            rel_part = related_parts[rId]
                                            print(f"        Resolved related part: {rel_part}")
                                            if hasattr(rel_part, 'target_ref'):
                                                print(f"        Related part target_ref: {rel_part.target_ref()}")
                                except Exception as e:
                                    print(f"        Error resolving: {e}")
        
        # Check shape-level hyperlinks (skip group shapes)
        is_group = hasattr(shape, "shapes")
        if not is_group:
            try:
                if hasattr(shape, "click_action"):
                    click_action = shape.click_action
                    print(f"  Has click_action: {click_action}")
                    if click_action:
                        print(f"    Click action type: {type(click_action)}")
                        print(f"    Click action dir: {dir(click_action)}")
            except Exception as e:
                if "group shape" not in str(e).lower():
                    print(f"  Click action error: {e}")

