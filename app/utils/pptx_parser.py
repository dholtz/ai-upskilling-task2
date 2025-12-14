"""Utility to parse PowerPoint files and extract text with URLs"""
from pptx import Presentation
from pptx.shapes.base import BaseShape
from typing import List, Dict
import logging

logger = logging.getLogger(__name__)

def extract_text_and_urls(pptx_path: str) -> List[Dict]:
    """
    Extract text and hyperlinks from a PowerPoint file.
    
    Args:
        pptx_path: Path to the .pptx file
        
    Returns:
        List of dictionaries containing slide data with text and URLs
    """
    try:
        prs = Presentation(pptx_path)
        slides_data = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []
            slide_urls = []
            seen_urls = set()  # Track URLs to avoid duplicates
            
            # Extract text and hyperlinks from all shapes
            for shape in slide.shapes:
                # Handle tables separately (they store hyperlinks differently)
                # Tables are stored as GraphicFrame objects with a .table attribute
                if hasattr(shape, 'table'):
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            # Get cell text
                            if cell.text and cell.text.strip():
                                slide_text_parts.append(cell.text.strip())
                            
                            # Extract hyperlinks from table cell
                            if hasattr(cell, "text_frame"):
                                for paragraph in cell.text_frame.paragraphs:
                                    cell_para_text = ""
                                    for run in paragraph.runs:
                                        if run.text:
                                            cell_para_text += run.text
                                        
                                        # Check for hyperlinks in table cell runs
                                        if hasattr(run, "hyperlink") and run.hyperlink:
                                            try:
                                                hyperlink = run.hyperlink
                                                url = getattr(hyperlink, 'address', None)
                                                
                                                # If no address, try to resolve through XML/rId
                                                if not url and hasattr(hyperlink, '_element'):
                                                    try:
                                                        element = hyperlink._element
                                                        rId = element.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') or element.get('rId')
                                                        if rId and hasattr(slide, 'part') and hasattr(slide.part, 'rels'):
                                                            rels = slide.part.rels
                                                            if rId in rels:
                                                                rel = rels[rId]
                                                                if hasattr(rel, 'target_ref'):
                                                                    url = rel.target_ref
                                                    except Exception:
                                                        pass
                                                
                                                if url:
                                                    url = str(url).strip()
                                                    if url and url not in seen_urls and not url.startswith('#'):
                                                        seen_urls.add(url)
                                                        link_text = run.text.strip() if run.text else cell_para_text.strip()
                                                        slide_urls.append({
                                                            'url': url,
                                                            'text': link_text[:200] if link_text else '',
                                                            'slide': slide_num
                                                        })
                                            except Exception as e:
                                                logger.debug(f"Error extracting hyperlink from table cell: {e}")
                                                continue
                    continue  # Skip to next shape after processing table
                
                # Get text from shape
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        slide_text_parts.append(text)
                
                # Extract hyperlinks from text runs and paragraphs
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = ""
                        paragraph_urls = []
                        
                        # First, collect all text from runs in the paragraph
                        for run in paragraph.runs:
                            if run.text:
                                paragraph_text += run.text
                        
                        # Check hyperlinks at multiple levels
                        # Method 1: Check each run for hyperlinks
                        for run in paragraph.runs:
                            if hasattr(run, "hyperlink"):
                                try:
                                    hyperlink = run.hyperlink
                                    url = None
                                    
                                    # Get URL from hyperlink
                                    if hyperlink:
                                        # Try direct address attribute first
                                        url = getattr(hyperlink, 'address', None)
                                        
                                        # If address is None, try to get rId from XML element and resolve it
                                        if not url and hasattr(hyperlink, '_element'):
                                            try:
                                                element = hyperlink._element
                                                # Check for rId in XML attributes
                                                rId = element.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') or element.get('rId')
                                                
                                                if rId:
                                                    # Try to resolve through the slide's part relationships
                                                    if hasattr(slide, 'part') and hasattr(slide.part, 'rels'):
                                                        rels = slide.part.rels
                                                        if rId in rels:
                                                            rel = rels[rId]
                                                            # Get the target reference
                                                            if hasattr(rel, 'target_ref'):
                                                                url = rel.target_ref
                                                            elif hasattr(rel, 'target_uri'):
                                                                url = rel.target_uri
                                                            elif hasattr(rel, '_target'):
                                                                url = rel._target
                                            except Exception as xml_error:
                                                logger.debug(f"Could not get URL from XML/rId: {xml_error}")
                                                pass
                                        
                                        # Also try the rId attribute if it exists
                                        if not url and hasattr(hyperlink, "rId"):
                                            rId = hyperlink.rId
                                            try:
                                                # Try to resolve through the slide's part
                                                if hasattr(slide, 'part') and hasattr(slide.part, 'rels'):
                                                    rels = slide.part.rels
                                                    if rId in rels:
                                                        rel = rels[rId]
                                                        if hasattr(rel, 'target_ref'):
                                                            url = rel.target_ref
                                                        elif hasattr(rel, 'target_uri'):
                                                            url = rel.target_uri
                                            except Exception as resolve_error:
                                                logger.debug(f"Could not resolve rId {rId}: {resolve_error}")
                                                pass
                                    
                                    if url:
                                        url = str(url).strip()
                                        # Filter out empty URLs and internal links
                                        if url and url not in seen_urls and not url.startswith('#'):
                                            seen_urls.add(url)
                                            link_text = run.text.strip() if run.text else paragraph_text.strip()
                                            paragraph_urls.append({
                                                'url': url,
                                                'text': link_text[:200] if link_text else '',
                                                'slide': slide_num
                                            })
                                except Exception as e:
                                    logger.debug(f"Error extracting hyperlink from run: {e}")
                                    continue
                        
                        # Method 2: Check paragraph-level hyperlinks (less common but possible)
                        if hasattr(paragraph, "hyperlink"):
                            try:
                                hyperlink = paragraph.hyperlink
                                if hyperlink:
                                    url = getattr(hyperlink, 'address', None)
                                    if url:
                                        url = str(url).strip()
                                        if url and url not in seen_urls and not url.startswith('#'):
                                            seen_urls.add(url)
                                            paragraph_urls.append({
                                                'url': url,
                                                'text': paragraph_text.strip()[:200] if paragraph_text else '',
                                                'slide': slide_num
                                            })
                            except Exception as e:
                                logger.debug(f"Error extracting paragraph hyperlink: {e}")
                                pass
                        
                        # Add all URLs found in this paragraph
                        slide_urls.extend(paragraph_urls)
                
                # Also check if the shape itself has a hyperlink (for images, shapes, etc.)
                # Skip group shapes as they don't support click actions
                # Check if it's a group shape first
                is_group = hasattr(shape, "shapes")  # Group shapes have a shapes collection
                if not is_group and hasattr(shape, "click_action"):
                    try:
                        click_action = shape.click_action
                        if click_action and hasattr(click_action, "hyperlinks"):
                            for hyperlink in click_action.hyperlinks:
                                if hasattr(hyperlink, "address") and hyperlink.address:
                                    url = hyperlink.address.strip()
                                    if url and url not in seen_urls:
                                        seen_urls.add(url)
                                        # Try to get associated text
                                        link_text = ""
                                        if hasattr(shape, "text") and shape.text:
                                            link_text = shape.text.strip()
                                        slide_urls.append({
                                            'url': url,
                                            'text': link_text[:200] if link_text else '',
                                            'slide': slide_num
                                        })
                    except Exception as e:
                        # Silently skip shapes that don't support click actions (like groups)
                        # The error message "a group shape cannot have a click action" is expected
                        if "group shape" not in str(e).lower():
                            logger.debug(f"Skipping click_action check for shape: {type(shape).__name__} - {e}")
                        pass
            
            # Combine all text
            full_text = ' '.join(slide_text_parts)
            
            if full_text or slide_urls:
                slides_data.append({
                    'slide_number': slide_num,
                    'text': full_text,
                    'urls': slide_urls,
                    'url_count': len(slide_urls)
                })
        
        total_urls = sum(len(slide['urls']) for slide in slides_data)
        logger.info(f"Extracted {len(slides_data)} slides with {total_urls} URLs from {pptx_path}")
        return slides_data
        
    except Exception as e:
        logger.error(f"Error parsing PowerPoint file: {str(e)}")
        raise

